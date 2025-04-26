import streamlit as st
import PyPDF2
import io
import os
from PIL import Image
import pytesseract
from pptx import Presentation

def display_job_description(job_details):
    """Display job description and requirements in a structured format"""
    # This function is now handled in app.py with HTML/CSS formatting
    pass

def validate_file(uploaded_file):
    """Validate if the uploaded file is a valid PDF, PPT, or image file"""
    file_type = uploaded_file.type
    try:
        if file_type == "application/pdf":
            PyPDF2.PdfReader(io.BytesIO(uploaded_file.getvalue()))
            return True, "pdf"
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            Presentation(io.BytesIO(uploaded_file.getvalue()))
            return True, "pptx"
        elif file_type in ["image/jpeg", "image/jpg", "image/png"]:
            Image.open(io.BytesIO(uploaded_file.getvalue()))
            return True, "image"
        else:
            return False, None
    except Exception as e:
        st.error(f"파일 검증 중 오류 발생: {str(e)}")
        return False, None

def extract_text_from_file(file_path, file_type):
    """Extract text content from a file based on its type"""
    try:
        text = ""
        
        if file_type == "pdf":
            # Extract text from PDF
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text()
        
        elif file_type == "pptx":
            # Extract text from PowerPoint
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        
        elif file_type == "image":
            # Extract text from image using OCR with Korean language support
            try:
                img = Image.open(file_path)
                text = pytesseract.image_to_string(img, lang='kor+eng')
            except:
                # If pytesseract is not available or fails, provide a basic description
                text = "이미지 기반 이력서. 내용 추출을 위해 고급 OCR 기능이 필요합니다."
        
        # Translate university names from English to Korean for better matching
        text = translate_university_names(text)
        
        return text
    except Exception as e:
        st.error(f"파일에서 텍스트 추출 중 오류 발생: {str(e)}")
        return None

def translate_university_names(text):
    """Translate English university names to Korean for better matching in analysis"""
    # University name translations (English to Korean)
    university_translations = {
        # Top tier universities
        "Seoul National University": "서울대",
        "Yonsei University": "연세대",
        "Korea University": "고려대", 
        "KAIST": "카이스트",
        "Korea Advanced Institute of Science and Technology": "카이스트",
        
        # Good universities
        "Sungkyunkwan University": "성균관대",
        "Hanyang University": "한양대", 
        "Ewha Womans University": "이화여대",
        "Sogang University": "서강대",
        "Chung-Ang University": "중앙대",
        "Kyung Hee University": "경희대",
        "Hankuk University of Foreign Studies": "한국외대",
        
        # Other notable universities
        "POSTECH": "포항공대",
        "Pohang University of Science and Technology": "포항공대",
        "Inha University": "인하대",
        "Korea University of Technology and Education": "한국기술교육대",
        "Kookmin University": "국민대",
        "Konkuk University": "건국대",
        "Sejong University": "세종대",
        "Dongguk University": "동국대",
        "Hongik University": "홍익대"
    }
    
    # Degree translations
    degree_translations = {
        "Bachelor": "학사",
        "Bachelor's": "학사",
        "Bachelor of Science": "이학사",
        "Bachelor of Arts": "문학사",
        "Bachelor of Engineering": "공학사",
        "Master": "석사",
        "Master's": "석사",
        "Master of Science": "이학석사",
        "Master of Arts": "문학석사",
        "Master of Engineering": "공학석사",
        "PhD": "박사",
        "Ph.D.": "박사",
        "Doctor of Philosophy": "박사",
        "Doctorate": "박사"
    }
    
    # Technical skill translations
    skill_translations = {
        "Machine Learning": "머신러닝",
        "Deep Learning": "딥러닝",
        "Natural Language Processing": "자연어처리",
        "NLP": "자연어처리",
        "Computer Vision": "컴퓨터 비전",
        "Data Science": "데이터 사이언스",
        "Data Analysis": "데이터 분석",
        "Artificial Intelligence": "인공지능",
        "AI": "인공지능",
        "Software Development": "소프트웨어 개발",
        "Web Development": "웹 개발",
        "Mobile Development": "모바일 개발",
        "Full Stack": "풀스택",
        "Frontend": "프론트엔드",
        "Backend": "백엔드",
        "DevOps": "데브옵스",
        "Cloud Computing": "클라우드 컴퓨팅",
        "Database": "데이터베이스",
        "UI/UX": "UI/UX 디자인",
        "Project Management": "프로젝트 관리"
    }
    
    # Translate university names
    for eng_name, kor_name in university_translations.items():
        text = text.replace(eng_name, f"{eng_name} {kor_name}")
    
    # Translate degree names
    for eng_degree, kor_degree in degree_translations.items():
        text = text.replace(eng_degree, f"{eng_degree} {kor_degree}")
    
    # Translate technical skills
    for eng_skill, kor_skill in skill_translations.items():
        text = text.replace(eng_skill, f"{eng_skill} {kor_skill}")
    
    return text

def format_progress_bar(percentage):
    """Format a percentage into a progress bar for display"""
    if percentage < 0:
        percentage = 0
    if percentage > 100:
        percentage = 100
    
    return percentage / 100

def create_download_link(content, filename, text):
    """Create a download link for text content"""
    import base64
    encoded = base64.b64encode(content.encode()).decode()
    href = f'<a href="data:file/txt;base64,{encoded}" download="{filename}" class="download-btn">{text}</a>'
    return href

def test_translation(test_text):
    """
    Test function to check how university and keyword translation works
    
    Args:
        test_text (str): Text to test translation on
        
    Returns:
        str: Translated text with highlighted translations
    """
    original_text = test_text
    translated_text = translate_university_names(test_text)
    
    print("===== Original Text =====")
    print(original_text)
    print("\n===== Translated Text =====")
    print(translated_text)
    
    return translated_text
